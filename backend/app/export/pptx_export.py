"""Export the generated PPT outline as an actual .pptx file using python-pptx."""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def export_paper_to_pptx(paper, output_path: Path) -> Path:
    outline = json.loads(paper.ppt_outline_json or "[]")
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    accent = RGBColor(0x43, 0x38, 0xCA)

    for i, slide_data in enumerate(outline):
        layout = title_layout if i == 0 else bullet_layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", f"Slide {i + 1}")
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent

        bullets = slide_data.get("bullets", [])
        if i == 0:
            if bullets and slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = bullets[0]
            continue

        body = slide.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
        tf = body.text_frame
        tf.clear()
        if not bullets:
            tf.text = "See paper for details."
        else:
            tf.text = str(bullets[0])
            tf.paragraphs[0].font.size = Pt(18)
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = str(b)
                p.font.size = Pt(18)
                p.level = 0

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
